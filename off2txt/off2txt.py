#!/usr/bin/env python
# -*- coding: utf-8 -*-
# off2txt: extract ASCII/Unicode text from Office files to separate files
# https://github.com/bnomis/off2txt
# (c) Simon Blanchard

import codecs
import os
import os.path
import sys

from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from .logger import error, init_logging


def writerr(options, line, exception=None, set_exit_status=True):
    if set_exit_status:
        options.exit_status = 'error'
    options.stderr.write(line + '\n')
    if exception:
        error(line, exc_info=True)


def writerr_file_access(options, line, exception=None):
    if not options.suppress_file_access_errors:
        writerr(options, line, exception=exception, set_exit_status=False)


def check_file_access(options, path):
    # check file exists
    # will return true even for broken symlinks
    if not os.path.lexists(path):
        writerr_file_access(options, 'File does not exist: %s' % path)
        return False

    # double check for broken symlinks
    if os.path.islink(path):
        if not os.path.exists(path):
            writerr_file_access(options, 'Broken symlink: %s' % path)
            return False

    # check can open for read
    if os.path.isdir(path):
        if not os.access(path, os.R_OK):
            writerr_file_access(options, 'Directory is not readable: %s' % path)
            return False
    else:
        try:
            fp = open(path)
        except Exception as e:
            writerr_file_access(options, 'File is not readable: %s' % path)
            error('check_file_access: exception for %s: %s' % (path, e), exc_info=True)
            return False
        else:
            fp.close()

    return True


def powerpoint(options, filename):
    text_runs = []
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                stripped = paragraph.text.strip()
                if stripped:
                    text_runs.append(paragraph.text)
    return text_runs


def word(options, filename):
    text_runs = []
    doc = Document(filename)
    for paragraph in doc.paragraphs:
        stripped = paragraph.text.strip()
        if stripped:
            text_runs.append(paragraph.text)
    return text_runs


def excel(options, filename):
    text_runs = []
    wb = load_workbook(filename=filename)
    for s in wb.get_sheet_names():
        ws = wb.get_sheet_by_name(s)
        for row in ws.rows:
            this_row = []
            for cell in row:
                this_row.append(cell.value)
            text_runs.append(this_row)
    return text_runs


def write_line_txt(options, fp, line):
    fp.write(line)
    fp.write(u'\n\n')


def write_line_csv(options, fp, line):
    fp.write(u', '.join(line))
    fp.write(u'\n')


def write_text_runs(options, filename, text_runs, is_csv=False):
    if not text_runs:
        return

    options.did_extract = True

    # setup mode
    mode = 'a'
    if options.overwrite_output_files:
        mode = 'w'

    # setup line output
    if is_csv:
        write_line = write_line_csv
    else:
        write_line = write_line_txt

    try:
        with codecs.open(filename, mode, 'utf8') as fp:
            for l in text_runs:
                write_line(options, fp, l)
    except Exception as e:
        writerr(options, 'Exception writing output file: %s' % filename, exception=e)


def split_end_of_line(outlines_ascii, outlines_unicode, line_ascii, line_unicode):
    if line_ascii:
        val = u''.join(line_ascii)
        stripped = val.strip()
        if stripped:
            outlines_ascii.append(val)
    if line_unicode:
        val = u''.join(line_unicode)
        stripped = val.strip()
        if stripped:
            outlines_unicode.append(val)

def split_ascii_unicode(options, text_runs):
    if not text_runs:
        return None, None

    outlines_ascii = []
    outlines_unicode = []

    for l in text_runs:
        line_ascii = []
        line_unicode = []
        for c in l:
            val = ord(c)
            if val < 128:
                # end of line
                if c in ('\n', '\r'):
                    split_end_of_line(outlines_ascii, outlines_unicode, line_ascii, line_unicode)
                    line_ascii = []
                    line_unicode = []
                else:
                    line_ascii.append(c)
            else:
                line_unicode.append(c)
        split_end_of_line(outlines_ascii, outlines_unicode, line_ascii, line_unicode)

    return outlines_ascii, outlines_unicode


def is_ascii_cell(options, cell):
    for c in cell:
        val = ord(c)
        if val >= 128:
            return False
    return True


def split_ascii_unicode_csv(options, text_runs):
    if not text_runs:
        return None, None

    outlines_ascii = []
    outlines_unicode = []

    for line in text_runs:
        line_ascii = []
        line_unicode = []
        for cell in line:
            if is_ascii_cell(options, cell):
                line_ascii.append(cell)
                line_unicode.append('')
            else:
                line_unicode.append(cell)
                line_ascii.append('')

        outlines_ascii.append(line_ascii)
        outlines_unicode.append(line_unicode)

    return outlines_ascii, outlines_unicode


def output_filename_split(options, filename):
    ascii_name = output_filename(options, filename)
    unicode_name = output_filename(options, filename, is_unicode=True)
    return ascii_name, unicode_name


def make_outdir(options):
    if options.directory:
        if not os.path.exists(options.directory):
            path = ''
            for p in options.directory.split('/'):
                path = os.path.join(path, p)
                if not os.path.exists(path):
                    try:
                        os.mkdir(path)
                    except Exception as e:
                        writerr(options, 'Exception making directory: %s' % path, exception=e)


def output_filename(options, filename, is_unicode=False):
    if options.output:
        return options.output

    base, ext = os.path.splitext(filename)
    if options.directory:
        make_outdir(options)
        base = os.path.basename(base)
        base = os.path.join(options.directory, base)
    if options.split:
        if is_unicode:
            base = base + '-' + options.unicode
        else:
            base = base + '-' + options.ascii

    # catch for Excel files
    out_ext = options.extension
    if ext == '.xlsx' and out_ext == 'txt':
        out_ext = 'csv'
    return base + '.' + out_ext


def write_out_split(options, filename, text_runs):
    if not text_runs:
        return

    outlines_ascii, outlines_unicode = split_ascii_unicode(options, text_runs)
    ascii_name, unicode_name = output_filename_split(options, filename)

    write_text_runs(options, ascii_name, outlines_ascii)
    write_text_runs(options, unicode_name, outlines_unicode)


def write_out(options, filename, text_runs):
    if not text_runs:
        return

    if options.split:
        write_out_split(options, filename, text_runs)
    else:
        write_text_runs(options, output_filename(options, filename), text_runs)


def write_csv_split(options, filename, text_runs):
    if not text_runs:
        return

    outlines_ascii, outlines_unicode = split_ascii_unicode_csv(options, text_runs)
    ascii_name, unicode_name = output_filename_split(options, filename)

    write_text_runs(options, ascii_name, outlines_ascii, is_csv=True)
    write_text_runs(options, unicode_name, outlines_unicode, is_csv=True)


def write_csv(options, filename):
    text_runs = excel(options, filename)
    if text_runs:
        if options.split:
            write_csv_split(options, filename, text_runs)
        else:
            write_text_runs(options, output_filename(options, filename), text_runs, is_csv=True)


def off2txt_file(options, path):
    if not check_file_access(options, path):
        return

    base, ext = os.path.splitext(path)
    ext_to_proc = {
        '.docx': word,
        '.pptx': powerpoint,
        '.xlsx': excel
    }
    if ext not in ext_to_proc:
        writerr('Unknown extension: %s' % ext)
        return

    if ext in ('.docx', '.pptx'):
        write_out(options, path, ext_to_proc[ext](options, path))
    else:
        write_csv(options, path)


def off2txt(options):
    for f in options.files:
        off2txt_file(options, f)


def main(argv, stdin=None, stdout=None, stderr=None):
    from .options import parse_opts
    exit_statuses = {
        'extracted': 0,
        'no-extract': 1,
        'error': 2,
        'not-set': -1
    }

    options = parse_opts(argv, stdin=stdin, stdout=stdout, stderr=stderr)
    if not options:
        return exit_statuses['error']

    init_logging(options)

    # do the extractions
    try:
        off2txt(options)
    except KeyboardInterrupt:
        writerr(options, '\nInterrupted')
    except Exception as e:
        writerr(options, 'off2txt exception', exception=e)
        options.exit_status = 'error'

    if options.exit_status == 'not-set':
        if options.did_extract:
            options.exit_status = 'extracted'
        else:
            options.exit_status = 'no-extract'

    return exit_statuses[options.exit_status]


def run():
    sys.exit(main(sys.argv[1:]))


if __name__ == '__main__':
    run()


